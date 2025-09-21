from flask import Flask,flash,jsonify,session,request, send_file, render_template, send_from_directory
import os
import re
import subprocess
import pdfplumber
from werkzeug.utils import secure_filename
from fpdf import FPDF
import time
import random
import json
import logging
import itertools
from datetime import datetime  
import subprocess
from pptx import Presentation
from pptx.util import Inches
import base64
from io import BytesIO
import io
import sys
import asyncio
from weasyprint import HTML
import tempfile
import pdfkit
import random

from playwright.sync_api import sync_playwright

import os
import firebase_admin
import shutil
from firebase_admin import credentials,auth
from functools import wraps
from flask import redirect, url_for
from flask import render_template, make_response
from fastapi import FastAPI, Request
from aiogram import Bot, Dispatcher
from aiogram.types import Update
from aiogram.filters import Command
from aiogram.types import Message
import telebot


cred = credentials.Certificate("./abogida-b1c87-firebase-adminsdk-fbsvc-b7d80b3e83.json")
firebase_admin.initialize_app(cred)
 
app = Flask(__name__)
app.secret_key = 'your-secret-key'
UPLOAD_FOLDER = 'uploads'
PROCESSED_FOLDER = 'processed'
DATA_STORE_FOLDER = 'data_store'
STORAGE_FILE = 'file_store.json'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PROCESSED_FOLDER, exist_ok=True)
os.makedirs(DATA_STORE_FOLDER, exist_ok=True)
logging.basicConfig(level=logging.INFO)

pdf = FPDF()
pdf.add_page()

pdf.add_font("AbyssinicaSIL-R", "", "fonts/AbyssinicaSIL-R.ttf", uni=True)
pdf.set_font("AbyssinicaSIL-R", size=14) 
# Ethiopian date approximation
from datetime import datetime, timedelta

from datetime import datetime
#pinting five in one


# Manual Ethiopian date mapping per Gregorian date
MANUAL_DATE_MAP = {
"2025/Aug/06": "2017/11/30",
"2033/Aug/04": "2025/11/28",

"2025/Aug/07": "2017/12/01",
"2033/Aug/05": "2025/11/29",

"2025/Aug/08": "2017/12/02",
"2033/Aug/06": "2025/11/30",

"2025/Aug/09": "2017/12/03",
"2033/Aug/07": "2025/12/01",

"2025/Aug/10": "2017/12/04",
"2033/Aug/08": "2025/12/02",

"2025/Aug/11": "2017/12/05",
"2033/Aug/09": "2025/12/03",

"2025/Aug/12": "2017/12/06",
"2033/Aug/10": "2025/12/04",

"2025/Aug/13": "2017/12/07",
"2033/Aug/11": "2025/12/05",

"2025/Aug/14": "2017/12/08",
"2033/Aug/12": "2025/12/06",

"2025/Aug/15": "2017/12/09",
"2033/Aug/13": "2025/12/07",

"2025/Aug/16": "2017/12/10",
"2033/Aug/14": "2025/12/08",

"2025/Aug/17": "2017/12/11",
"2033/Aug/15": "2025/12/09",

"2025/Aug/18": "2017/12/12",
"2033/Aug/16": "2025/12/10",

"2025/Aug/19": "2017/12/13",
"2033/Aug/17": "2025/12/11",

"2025/Aug/20": "2017/12/14",
"2033/Aug/18": "2025/12/12",

"2025/Aug/21": "2017/12/15",
"2033/Aug/19": "2025/12/13",

"2025/Aug/22": "2017/12/16",
"2033/Aug/20": "2025/12/14",

"2025/Aug/23": "2017/12/17",
"2033/Aug/21": "2025/12/15",

"2025/Aug/24": "2017/12/18",
"2033/Aug/22": "2025/12/16",

"2025/Aug/25": "2017/12/19",
"2033/Aug/23": "2025/12/17",

"2025/Aug/26": "2017/12/20",
"2033/Aug/24": "2025/12/18",

"2025/Aug/27": "2017/12/21",
"2033/Aug/25": "2025/12/19",

"2025/Aug/28": "2017/12/22",
"2033/Aug/26": "2025/12/20",

"2025/Aug/29": "2017/12/23",
"2033/Aug/27": "2025/12/21",

"2025/Aug/30": "2017/12/24",
"2033/Aug/28": "2025/12/22",

"2025/Aug/31": "2017/12/25",
"2033/Aug/29": "2025/12/23",

"2025/Sep/01": "2017/12/26",
"2033/Aug/30": "2025/12/24",

"2025/Sep/02": "2017/12/27",
"2033/Aug/31": "2025/12/25",

"2025/Sep/03": "2017/12/28",
"2033/Sep/01": "2025/12/26",

"2025/Sep/04": "2017/12/29",
"2033/Sep/02": "2025/12/27",

"2025/Sep/05": "2017/12/30",
"2033/Sep/03": "2025/12/28",

"2025/Sep/06": "2017/13/01",
"2033/Sep/04": "2025/12/29",

"2025/Sep/07": "2017/13/02",
"2033/Sep/05": "2025/12/30",

"2025/Sep/08": "2017/13/03",
"2033/Sep/06": "2025/13/01",

"2025/Sep/09": "2017/13/04",
"2033/Sep/07": "2025/13/02",

"2025/Sep/10": "2017/13/05",
"2033/Sep/08": "2025/13/03",

"2025/Sep/11": "2018/01/01",
"2033/Sep/09": "2025/13/04",

"2025/Sep/12": "2018/01/02",
"2033/Sep/10": "2025/13/05",

"2025/Sep/13": "2018/01/03",
"2033/Sep/11": "2025/01/01",

"2025/Sep/14": "2018/01/04",
"2033/Sep/12": "2025/01/02",

"2025/Sep/15": "2018/01/05",
"2033/Sep/13": "2025/01/03",

"2025/Sep/16": "2018/01/06",
"2033/Sep/14": "2025/01/04",

"2025/Sep/17": "2018/01/07",
"2033/Sep/15": "2025/01/05",

"2025/Sep/18": "2018/01/08",
"2033/Sep/16": "2025/01/06",

"2025/Sep/19": "2018/01/09",
"2033/Sep/17": "2025/01/07",

"2025/Sep/20": "2018/01/10",
"2033/Sep/18": "2025/01/08"
}

def get_current_and_expiry_dates():
    today = datetime.now()
    english_date = today.strftime("%Y/%b/%d")

    # Get Ethiopian date manually
    if english_date not in MANUAL_DATE_MAP:
        raise ValueError(f"No Ethiopian mapping for {english_date}. Add it to MANUAL_DATE_MAP.")
    ethiopian_date = MANUAL_DATE_MAP[english_date]

    # Calculate expiry date (8 years ahead, 2 days earlier)
    expiry_date = datetime(today.year + 8, today.month, max(today.day - 2, 1))
    english_expiry = expiry_date.strftime("%Y/%b/%d")

    if english_expiry not in MANUAL_DATE_MAP:
        raise ValueError(f"No Ethiopian expiry mapping for {english_expiry}. Add it to MANUAL_DATE_MAP.")
    ethiopian_expiry = MANUAL_DATE_MAP[english_expiry]

    return english_date, ethiopian_date, ethiopian_expiry, english_expiry


def save_processed_path(original_filename, processed_path):
    data = {}
    if os.path.exists(STORAGE_FILE):
        with open(STORAGE_FILE, 'r') as f:
            data = json.load(f)
    data[original_filename] = processed_path
    with open(STORAGE_FILE, 'w') as f:
        json.dump(data, f)

def get_processed_path(filename):
    if not os.path.exists(STORAGE_FILE):
        return None
    with open(STORAGE_FILE, 'r') as f:
        data = json.load(f)
        return data.get(filename)

def extract_info_from_text(text,fan):
    lines = text.strip().split("\n")
    try:
        return {
            'Fin': lines[0],
            'amharic_name': lines[1],
            'english_name': lines[2],
            'ethiopian_bd': lines[3],
            'european_bd': lines[4].split()[0],
            'amharic_city': lines[4].split()[1],
            'english_city': lines[5],
            'amharic_gender': lines[6],
            'english_gender': lines[7].split()[0],
            'subcity': lines[7].split()[1],
            'english_subcity': lines[8],
            'amharic_nationality': lines[9],
            'english_nationality': lines[10],
            'amharic_woreda': lines[11],
            'english_woreda': lines[12],
            'phone': lines[13],
            "fan": fan
        }
    except Exception as e:
        logging.error(f"Error parsing lines: {e}")
        return None

def generate_id_pdf(info_list):
    pdf = FPDF("P", "mm", "A4")
    pdf.add_page()
    pdf.set_font("Helvetica", size=10)
    x_start, y_start = 10, 0
    card_width, card_height = 90, 50
    spacing_x, spacing_y = 5, 10

    for i, info in enumerate(info_list):
        x = x_start + (i % 2) * (card_width + spacing_x)
        y = y_start + (i // 2) * (card_height + spacing_y)





        # Add the Abyssinica font from your local font file, you can give it a name
        pdf.add_font('AbyssinicaSIL-R', '', 'fonts/AbyssinicaSIL-R.ttf', uni=True)

        pdf.add_page()
        pdf.set_font("AbyssinicaSIL-R", size=14)

        pdf.cell(0, 10, "Some text in Abyssinica font", ln=True)



        pdf.set_font("AbyssinicaSIL-R", size=14) 

        pdf.set_xy(x, y)
        pdf.cell(0, 5, f"Name: {info['english_name']}", ln=1)
        pdf.set_x(x)
        pdf.cell(0, 5, f"Gender: {info['english_gender']}", ln=1)
        pdf.set_x(x)
        pdf.cell(0, 5, f"DOB: {info['european_bd']}", ln=1)
        pdf.set_x(x)
        pdf.cell(0, 5, f"City: {info['english_city']}", ln=1)
        pdf.set_x(x)
        pdf.cell(0, 5, f"Phone: {info['phone']}", ln=1)

    output_path = os.path.join(PROCESSED_FOLDER, f"id_cards_{int(time.time())}.pdf")
    pdf.output(output_path)
    return output_path

@app.before_request
def add_zip_to_jinja():
    app.jinja_env.globals.update(zip=itertools.zip_longest)


def login_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function


@app.route('/verify_token', methods=['POST'])
def verify_token():
    data = request.get_json()
    id_token = data.get("idToken")
    try:
        decoded_token = auth.verify_id_token(id_token)
        session['user'] = {
            'uid': decoded_token['uid'],
            'email': decoded_token.get('email')
        }
        return jsonify({"success": True})
    except Exception as e:
        logging.error(f"Token verification failed: {e}")
        return jsonify({"success": False}), 401

@app.route('/login')
def login():
    return render_template('login.html')


@app.route('/upload', methods=['POST'])
def upload():
    if 'files' not in request.files:
        return "No files part in request", 400

    files = request.files.getlist('files')
    if not files or all(file.filename == '' for file in files):
        return "No files selected", 400

    output_paths = []

    for file in files:
        if file.filename == '':
            continue

        filename = secure_filename(file.filename)
        fan = os.path.splitext(filename)[0]  # → 'fan1234'
        fan = ''.join(re.findall(r'\d+', fan))[:12]
        fan = '-'.join([fan[i:i+4] for i in range(0, len(fan), 4)])
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)

        # Step 1: Run PDF preprocessing script
        try:
            subprocess.run(["python", "generate_data_pdf.py", filepath, fan], check=True)
        except subprocess.CalledProcessError as e:
            return f"Error processing {filename}: {e}", 500

        # Step 2: Get current and expiry dates
        english_date, ethiopian_date, english_expiry, ethiopian_expiry = get_current_and_expiry_dates()

        # Step 3: Extract data from PDF
        extracted_infos = []
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                info = extract_info_from_text(text, fan)
                if info:
                    info['generated_date_gregorian'] = english_date
                    info['generated_date_ethiopian'] = ethiopian_date
                    info['generated_expiry_date_gregorian'] = english_expiry
                    info['generated_expiry_date_ethiopian'] = ethiopian_expiry
                    extracted_infos.append(info)
                if len(extracted_infos) == 5:
                    break

        if not extracted_infos:
            continue  # Skip files with no extractable info

        # Step 4: Save JSON and generate output PDF
        json_path = os.path.join(DATA_STORE_FOLDER, filename + '.json')
    

        # Check if file exists and is not empty
    if not os.path.exists(json_path) or os.path.getsize(json_path) == 0:
                data = {}
    else:
     with open(json_path, "r", encoding="utf-8") as f:
                    try:
                        data = json.load(f)
                    except json.JSONDecodeError:
                        data = {}
            
    data[filename] = output_pdf_path
            


    with open(json_path, 'w', encoding='utf-8') as f:
                    json.dump(extracted_infos, f, ensure_ascii=False, indent=2)


    output_pdf_path = generate_id_pdf(extracted_infos)
    save_processed_path(filename, output_pdf_path)
    output_paths.append(output_pdf_path)

    if not output_paths:
        return "files processed successfully", 500

    return f"Processed {len(output_paths)} file(s) successfully."
@app.route('/get-processed/<filename>')
def get_processed(filename):
    path = get_processed_path(filename)
    if not path or not os.path.exists(path):
        return "Processed file not found", 404
    return send_file(path, as_attachment=True)

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))
@app.route('/print')
def print_ids():
    id_cards_data = []
    mugshot_images = []
    barcode_images = []
    qr_code_images = []
    english_date, ethiopian_date ,ethiopian_expiry, english_expiry = get_current_and_expiry_dates()

    for folder in os.listdir(DATA_STORE_FOLDER):
        folder_path = os.path.join(DATA_STORE_FOLDER, folder)
        if not os.path.isdir(folder_path):
            continue

        info_path = os.path.join(folder_path, 'info.json')
        mugshot_path = os.path.join(folder_path, 'photo.png')
        barcode_path = os.path.join(folder_path, 'barcode.png')
        qr_path = os.path.join(folder_path, 'qr.png')

        if not (os.path.exists(info_path) and os.path.exists(mugshot_path)
                and os.path.exists(barcode_path) and os.path.exists(qr_path)):
            continue

        with open(info_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        id_cards_data.append(data)
        mugshot_images.append(f"/data_store/{folder}/photo.png")
        barcode_images.append(f"/data_store/{folder}/barcode.png")
        qr_code_images.append(f"/data_store/{folder}/qr.png")
        person_ids = [int(folder.split("_")[1]) for folder in sorted(os.listdir("data_store")) if folder.startswith("person_")]
        person_ids.sort()  
    return render_template('print.html',
                           id_cards_data=id_cards_data,
                           mugshot_images=mugshot_images,
                           barcode_images=barcode_images,
                           qr_code_images=qr_code_images, 
                           english_date=english_date,
                           ethiopian_date=ethiopian_date,
                           english_expiry=english_expiry,
                           ethiopian_expiry=ethiopian_expiry,
                           person_ids=person_ids,
                           random=random)


@app.route('/print_color')
def print_color():
    id_cards_data = []
    mugshot_images = []
    barcode_images = []
    qr_code_images = []
    english_date, ethiopian_date ,ethiopian_expiry, english_expiry = get_current_and_expiry_dates()

    for folder in os.listdir(DATA_STORE_FOLDER):
        folder_path = os.path.join(DATA_STORE_FOLDER, folder)
        if not os.path.isdir(folder_path):
            continue

        info_path = os.path.join(folder_path, 'info.json')
        mugshot_path = os.path.join(folder_path, 'photo.png')
        barcode_path = os.path.join(folder_path, 'barcode.png')
        qr_path = os.path.join(folder_path, 'qr.png')

        if not (os.path.exists(info_path) and os.path.exists(mugshot_path)
                and os.path.exists(barcode_path) and os.path.exists(qr_path)):
            continue

        with open(info_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        id_cards_data.append(data)
        mugshot_images.append(f"/data_store/{folder}/photo.png")
        barcode_images.append(f"/data_store/{folder}/barcode.png")
        qr_code_images.append(f"/data_store/{folder}/qr.png")

    person_ids = [int(folder.split("_")[1]) for folder in sorted(os.listdir("data_store")) if folder.startswith("person_")]
    person_ids.sort()

    # 👉 you can use a separate template print_color.html or reuse print.html
    return render_template('print_color.html',
                           id_cards_data=id_cards_data,
                           mugshot_images=mugshot_images,
                           barcode_images=barcode_images,
                           qr_code_images=qr_code_images,
                           english_date=english_date,
                           ethiopian_date=ethiopian_date,
                           english_expiry=english_expiry,
                           ethiopian_expiry=ethiopian_expiry,
                           person_ids=person_ids,
                           random=random)

from flask import send_from_directory
@app.route('/data_store/<path:filename>')
def data_store_static(filename):
    return send_from_directory('data_store', filename)

@app.route('/data_store/<person_id>/<filename>')
# @login_required
def serve_data_store_images(person_id, filename):
    return send_from_directory(os.path.join(DATA_STORE_FOLDER, person_id), filename)

@app.route('/')
# @login_required
def index():
    return render_template('upload.html')

@app.route('/delete')
@login_required
def delete_page():
    id_cards = []
    for folder in os.listdir(DATA_STORE_FOLDER):
        folder_path = os.path.join(DATA_STORE_FOLDER, folder)
        if not os.path.isdir(folder_path):
            continue

        info_path = os.path.join(folder_path, 'info.json')
        mugshot_path = os.path.join(folder_path, 'photo.png')

        if not (os.path.exists(info_path) and os.path.exists(mugshot_path)):
            continue

        with open(info_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        data['photo'] = f"/data_store/{folder}/photo.png"
        data['person_id'] = folder  

        id_cards.append(data)

    return render_template('delete.html', id_cards_data=id_cards)

@app.route('/delete/<person_id>', methods=['POST'])
@login_required
def delete_person(person_id):
    folder_path = os.path.join(DATA_STORE_FOLDER, person_id)

    if os.path.exists(folder_path) and os.path.isdir(folder_path):
        try:
            import shutil
            shutil.rmtree(folder_path)
            return jsonify({"success": True, "message": f"{person_id} deleted."})
        except Exception as e:
            logging.error(f"Error deleting {person_id}: {e}")
            return jsonify({"success": False, "error": str(e)}), 500
    else:
        return jsonify({"success": False, "error": "Folder not found"}), 404



UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)



@app.route('/api/list', methods=['GET'])
def api_list_records():
    id_cards = []
    for folder in os.listdir(DATA_STORE_FOLDER):
        folder_path = os.path.join(DATA_STORE_FOLDER, folder)
        if not os.path.isdir(folder_path):
            continue

        info_path = os.path.join(folder_path, 'info.json')
        mugshot_path = os.path.join(folder_path, 'photo.png')

        if not (os.path.exists(info_path) and os.path.exists(mugshot_path)):
            continue

        with open(info_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        data['photo'] = f"/data_store/{folder}/photo.png"
        data['person_id'] = folder
        id_cards.append(data)

    return jsonify({"success": True, "records": id_cards})


@app.route('/delete_up', methods=['GET', 'POST'])
def delete_files():
    if request.method == 'POST':
        deleted_files = []
        for filename in os.listdir(UPLOAD_FOLDER):
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            if os.path.isfile(file_path) and filename.lower().endswith('.pdf'):
                os.remove(file_path)
                deleted_files.append(filename)

        flash(f"Deleted {len(deleted_files)} PDF(s) from uploads.", "success")
        return redirect(url_for('delete_files'))

    uploaded_files = [
        f for f in os.listdir(UPLOAD_FOLDER)
        if os.path.isfile(os.path.join(UPLOAD_FOLDER, f)) and f.lower().endswith('.pdf')
    ]
    return render_template('delete_uploads.html', files=uploaded_files)


@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    try:
        subprocess.run(["node", "screenshot_id_cards.js"], check=True)
        subprocess.run(["node", "export_to_ppt.js"], check=True)
        return {"success": True, "message": "PPT created successfully."}
    except subprocess.CalledProcessError as e:
        import traceback
        traceback.print_exc()  # ← Print detailed error to the terminal
        return {"success": False, "error": str(e)}, 500
@app.route('/delete_all', methods=['POST'])
def delete_all_folders():
    data_store_path = os.path.join(os.getcwd(), 'data_store')
    deleted_folders = []

    try:
        if os.path.exists(data_store_path):
            for item in os.listdir(data_store_path):
                item_path = os.path.join(data_store_path, item)
                if os.path.isdir(item_path):
                    shutil.rmtree(item_path)
                    deleted_folders.append(item)
        return jsonify({
            "success": True,
            "message": f"Deleted {len(deleted_folders)} folder(s) from data_store."
        })
    except Exception as e:
        return jsonify({
            "success": False,
            "error": str(e)
        })
@app.route('/api/upload', methods=['POST'])
def api_upload():
    if 'file' not in request.files:
        return jsonify({"success": False, "error": "No file uploaded"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"success": False, "error": "Empty filename"}), 400

    try:
        filename = secure_filename(file.filename)
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)

        # process the file (reuse your logic from upload())
        fan = os.path.splitext(filename)[0]
        fan = ''.join(re.findall(r'\d+', fan))[:12]
        fan = '-'.join([fan[i:i+4] for i in range(0, len(fan), 4)])

        # Call your processing pipeline
        subprocess.run(["python", "generate_data_pdf.py", filepath, fan], check=True)

        eng_date, eth_date, eng_exp, eth_exp = get_current_and_expiry_dates()
        output_pdf_path = generate_id_pdf([{
            "english_name": "Demo", 
            "english_gender": "M", 
            "european_bd": "2000-01-01", 
            "english_city": "Addis", 
            "phone": "0912345678"
        }])  # 👈 replace with real extraction logic

        save_processed_path(filename, output_pdf_path)

        return jsonify({
            "success": True,
            "message": f"{filename} processed successfully",
            "generated_pdf": output_pdf_path,
            "dates": {
                "english_date": eng_date,
                "ethiopian_date": eth_date,
                "expiry_english": eng_exp,
                "expiry_ethiopian": eth_exp
            }
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

# --- Route to generate PDF ---
async def generate_pdf(html_content):
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page()
        await page.set_content(html_content)
        pdf_bytes = await page.pdf(format='A4', print_background=True)
        await browser.close()
        return pdf_bytes

# Sample data
id_cards_data = [
    {
        "phone": "0912345678",
        "amharic_nationality": "ኢትዮጵያዊ",
        "english_nationality": "Ethiopian",
        "amharic_name": "ማላኩ ዘገዬ",
        "english_name": "Melaku Zegeye",
        "ethiopian_bd": "2000-01-01",
        "amharic_city": "አዲስ አበባ",
        "english_city": "Addis Ababa",
        "amharic_subcity": "ካምፓርት",
        "english_subcity": "Kampart",
        "amharic_woreda": "ወረዳ 1",
        "english_woreda": "Woreda 1",
        "amharic_gender": "ወንድ",
        "english_gender": "Male",
        "Fin": "FIN123456",
        "fan": "FAN98765"
    },
    # add more cards here
]

mugshot_images = ["/data_store/person_001/photo.png"]  # or url_for('static', filename='...')
barcode_images = ["/data_store/person_001/barcode.png"]
qr_code_images = ["/data_store/person_001/qr.png"]
person_ids = [1]

@app.route("/api/print_pdf")
def print_pdf():
    # Render HTML with variables
    rendered_html = render_template(
        "print.html",
        id_cards_data=id_cards_data,
        mugshot_images=mugshot_images,
        barcode_images=barcode_images,
        qr_code_images=qr_code_images,
        person_ids=person_ids,
        english_date="2025-08-30",
        ethiopian_date="2017-12-25",
        english_expiry="2030-08-30",
        ethiopian_expiry="2023-12-25",
        random=random 
    )

     # Save HTML to temp file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".html") as tmp_html:
        tmp_html.write(rendered_html.encode("utf-8"))
        tmp_html_path = tmp_html.name

    # Output PDF path
    pdf_path = tmp_html_path.replace(".html", ".pdf")

    # Convert HTML → PDF using wkhtmltopdf with enable-local-file-access
    options = {
        "enable-local-file-access": ""
    }
    pdfkit.from_file(tmp_html_path, pdf_path, options=options)

    # Return PDF as download
    return send_file(pdf_path, as_attachment=True, download_name="ethiopian-id-cards.pdf")










def emulate_delete():
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)  # headless browser
        page = browser.new_page()
        page.goto("http://localhost:7000/print")  # your print.html endpoint

        # Wait until the page loads fully
        page.wait_for_load_state("networkidle")

        # Click the Delete button by selector
        page.click("a[href='/delete']")  # assuming your button is <a href="/delete">

        # Optional: wait for deletion to finish or page to update
        page.wait_for_timeout(1000)

        browser.close()


@app.route("/trigger_delete")
def trigger_delete():
    emulate_delete()  # call the headless browser function
    return "Delete button triggered!"













def emulate_save_pdf():
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)  # headless browser
        page = browser.new_page()
        
        # Use 127.0.0.1 to avoid localhost resolution issues
        page.goto("http://127.0.0.1:7000/print", wait_until="networkidle", timeout=60000)

        # Wait a little if the page has heavy JS rendering
        page.wait_for_timeout(2000)  # 2 seconds

        # Click the Save as PDF button
        # Make sure your button has an onclick like: <button onclick="savePDF()">
        page.click("button:has-text('Save as PDF')")

        # Optional: wait until PDF generation finishes (depends on your JS)
        page.wait_for_timeout(5000)  # 5 seconds

        browser.close()


@app.route("/trigger_save_pdf")
def trigger_save_pdf():
    emulate_save_pdf()  # call the headless browser function
    return "pdf downloaded!"

@app.route('/generate_pptx_from_images', methods=['POST'])
def generate_pptx():
    data = request.get_json()
    images = data.get("images", [])

    # ✅ A4 slide size (in inches)
    ppt_width = Inches(7.5)        # ~190.5 mm (standard printable A4 width)
    ppt_height = Inches(10.833)    # ~275 mm

    # ✅ ID card size and layout settings
    id_width_inch = 7.03           # Width of each ID
    id_height_inch = 2.13          # Height of each ID
    spacing_inch = 0.045           # Vertical space between cards
    cards_per_slide = 5

    # Create presentation
    prs = Presentation()
    prs.slide_width = ppt_width
    prs.slide_height = ppt_height

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    for i, img_b64 in enumerate(images):
        try:
            img_data = base64.b64decode(img_b64.split(",")[1])
            img_stream = BytesIO(img_data)

            # 🆕 Add new slide after every 5 cards
            if i % cards_per_slide == 0 and i != 0:
                slide = prs.slides.add_slide(prs.slide_layouts[6])

            # ✅ FIXED: Ensure top is converted to Inches
            vertical_offset = (id_height_inch + spacing_inch) * (i % cards_per_slide)
            top = Inches(vertical_offset)
            left = Inches(0.25)  # optional horizontal padding

            # 🖼️ Add picture
            slide.shapes.add_picture(img_stream, left, top,
                                     width=Inches(id_width_inch),
                                     height=Inches(id_height_inch))

        except Exception as e:
            print(f"Error adding image {i}: {e}")

    # 📤 Return .pptx as a download
    output_stream = BytesIO()
    prs.save(output_stream)
    output_stream.seek(0)

    return send_file(
        output_stream,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name="ID_Cards.pptx"
    )

@app.context_processor
def inject_user():
 return dict(user=session.get('user'))






BOT_TOKEN = os.getenv("BOT_TOKEN")
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher()

# ✅ import your bot handlers
from bot import router
dp.include_router(router)

app = FastAPI()



# Set webhook on startup
@app.on_event("startup")
async def on_startup():
    webhook_url = f"{BASE_URL}/webhook/{BOT_TOKEN}"
    await bot.set_webhook(webhook_url)

# Handle Telegram updates
@app.post("/webhook/{token}")
async def telegram_webhook(token: str, request: Request):
    if token != BOT_TOKEN:
        return {"error": "Invalid token"}
    data = await request.json()
    update = Update(**data)
    await dp.feed_update(bot, update)
    return {"status": "ok"}

# Handlers
@dp.message(Command("start"))
async def start_handler(message: Message):
    await message.answer("👋 Hello! Bot is running on webhook mode.")

@dp.message(Command("upload"))
async def upload_handler(message: Message):
    await message.answer("📂 Please send me a PDF file to upload.")


if __name__ == '__main__':
    app.run(debug=True, port=7000)
